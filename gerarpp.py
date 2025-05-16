from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import os

def create_formatted_cv_presentation(data_df, output_path="Employee_Profiles.pptx"):
    prs = Presentation()

    # Slide 4:3 (10x7.5 inches) já definido assim:
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    PURPLE_COLOR = RGBColor(128, 0, 128)
    DARK_GRAY = RGBColor(64, 64, 64)
    BLUE_COLOR = RGBColor(0, 0, 255)

    for idx, employee in data_df.iterrows():
        slide = prs.slides.add_slide(blank_layout)

        # **Remover borda preta: comentar ou retirar este bloco**
        # border = slide.shapes.add_shape(
        #     MSO_SHAPE.RECTANGLE,
        #     Inches(0.2), Inches(0.2),
        #     Inches(9.6), Inches(7.1)
        # )
        # border.fill.background()  # Transparente
        # border.line.color.rgb = RGBColor(0, 0, 0)
        # border.line.width = Pt(1.0)

        # Placeholder círculo azul (igual)
        left = Inches(0.4)
        top = Inches(0.4)
        width = height = Inches(1)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, top, width, height
        )
        circle.fill.background()
        circle.line.color.rgb = BLUE_COLOR
        circle.line.width = Pt(1.5)

        # Texto do placeholder dentro do círculo
        pic_text = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(0.5))
        pic_text_frame = pic_text.text_frame
        pic_text_frame.word_wrap = True
        pic_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = pic_text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Place\nPicture here\nOr Delete"
        font = run.font
        font.name = "Arial"
        font.size = Pt(8)  # menor fonte aqui
        font.bold = True
        font.color.rgb = BLUE_COLOR

        # Nome - diminuir fonte de 36 para 28
        name_box = slide.shapes.add_textbox(Inches(2.2), Inches(0.4), Inches(6), Inches(0.8))
        name_frame = name_box.text_frame
        name_p = name_frame.paragraphs[0]
        name_p.alignment = PP_ALIGN.CENTER
        name_run = name_p.add_run()
        name_run.text = employee["First Name Last Name"]
        name_font = name_run.font
        name_font.name = "Arial"
        name_font.size = Pt(28)
        name_font.color.rgb = DARK_GRAY

        # Job title - diminuir para 14pt
        job_box = slide.shapes.add_textbox(Inches(2.2), Inches(1.1), Inches(6), Inches(0.5))
        job_frame = job_box.text_frame
        job_p = job_frame.paragraphs[0]
        job_p.alignment = PP_ALIGN.CENTER
        job_run = job_p.add_run()
        job_run.text = employee["Job Title / Role"]
        job_font = job_run.font
        job_font.name = "Arial"
        job_font.size = Pt(14)
        job_font.color.rgb = DARK_GRAY

        # Profile header - diminuir para 18pt
        profile_header = slide.shapes.add_textbox(Inches(0.4), Inches(1.8), Inches(2), Inches(0.4))
        profile_frame = profile_header.text_frame
        profile_p = profile_frame.paragraphs[0]
        profile_run = profile_p.add_run()
        profile_run.text = "Profile"
        profile_font = profile_run.font
        profile_font.name = "Arial"
        profile_font.size = Pt(18)
        profile_font.bold = True
        profile_font.color.rgb = PURPLE_COLOR

        # Profile content - diminuir para 9pt
        profile_content = slide.shapes.add_textbox(Inches(0.4), Inches(2.2), Inches(5), Inches(1.1))
        profile_content_frame = profile_content.text_frame
        profile_content_frame.word_wrap = True
        profile_content_p = profile_content_frame.paragraphs[0]
        profile_content_run = profile_content_p.add_run()
        profile_content_run.text = employee["Profile"]
        profile_content_font = profile_content_run.font
        profile_content_font.name = "Arial"
        profile_content_font.size = Pt(9)
        profile_content_font.color.rgb = DARK_GRAY

        # Education header - diminuir para 18pt
        edu_header = slide.shapes.add_textbox(Inches(0.4), Inches(2.8), Inches(2), Inches(0.4))
        edu_frame = edu_header.text_frame
        edu_p = edu_frame.paragraphs[0]
        edu_run = edu_p.add_run()
        edu_run.text = "Education"
        edu_font = edu_run.font
        edu_font.name = "Arial"
        edu_font.size = Pt(18)
        edu_font.bold = True
        edu_font.color.rgb = PURPLE_COLOR

        # Education content - diminuir para 9pt
        edu_content = slide.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(5), Inches(1.0))
        edu_content_frame = edu_content.text_frame
        edu_content_frame.word_wrap = True
        edu_content_p = edu_content_frame.paragraphs[0]
        edu_content_run = edu_content_p.add_run()
        edu_content_run.text = employee["Education"]
        edu_content_font = edu_content_run.font
        edu_content_font.name = "Arial"
        edu_content_font.size = Pt(9)
        edu_content_font.color.rgb = DARK_GRAY

        # Skills header - Arial 18pt, negrito, roxo
        skills_header = slide.shapes.add_textbox(Inches(0.4), Inches(3.8), Inches(5), Inches(0.4))
        skills_frame = skills_header.text_frame
        skills_p = skills_frame.paragraphs[0]
        skills_run = skills_p.add_run()
        skills_run.text = "Relevant Skills & Qualifications"
        skills_font = skills_run.font
        skills_font.name = "Arial"
        skills_font.size = Pt(18)
        skills_font.bold = True
        skills_font.color.rgb = PURPLE_COLOR

        # Ajuste das skills para 2 colunas para não sair do slide
        skills_list = employee["Relevant Skills & Qualifications"].split('\n')
        max_rows_per_col = 8  # para não ficar muito longo verticalmente

        # Ajustar número de linhas para a tabela
        rows_needed = min(len(skills_list), max_rows_per_col)

        # Coluna 1
        skills_table_left = slide.shapes.add_table(
            rows=rows_needed,
            cols=1,
            left=Inches(0.5),
            top=Inches(4.2),
            width=Inches(2.2),
            height=Inches(1.5)
        ).table

        # Coluna 2 (se tiver mais skills)
        skills_table_right = None
        if len(skills_list) > max_rows_per_col:
            skills_table_right = slide.shapes.add_table(
                rows=len(skills_list) - max_rows_per_col,
                cols=1,
                left=Inches(3.0),
                top=Inches(4.2),
                width=Inches(2.2),
                height=Inches(1.5)
            ).table

        # Preencher coluna 1
        for i in range(rows_needed):
            cell = skills_table_left.cell(i, 0)
            cell.text = skills_list[i]
            cell.fill.background()
            cell.text_frame.word_wrap = True
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = "Arial"
                paragraph.font.color.rgb = DARK_GRAY

        # Preencher coluna 2 se existir
        if skills_table_right:
            for i in range(len(skills_list) - max_rows_per_col):
                cell = skills_table_right.cell(i, 0)
                cell.text = skills_list[i + max_rows_per_col]
                cell.fill.background()
                cell.text_frame.word_wrap = True
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.name = "Arial"
                    paragraph.font.color.rgb = DARK_GRAY


        # Experience header - diminuir para 18pt
        exp_header = slide.shapes.add_textbox(Inches(6.0), Inches(1.8), Inches(3.5), Inches(0.4))
        exp_frame = exp_header.text_frame
        exp_p = exp_frame.paragraphs[0]
        exp_run = exp_p.add_run()
        exp_run.text = "Relevant Experience"
        exp_font = exp_run.font
        exp_font.name = "Arial"
        exp_font.size = Pt(18)
        exp_font.bold = True
        exp_font.color.rgb = PURPLE_COLOR

        # Experience content - diminuir fonte para 9pt e mover um pouco pra cima
        exp_content = slide.shapes.add_textbox(Inches(6.0), Inches(2.3), Inches(3.5), Inches(3.6))
        exp_content_frame = exp_content.text_frame
        exp_content_frame.word_wrap = True

        experience_lines = employee["Relevant Experience"].split('\n')
        if experience_lines:
            first_p = exp_content_frame.paragraphs[0]
            first_p.clear()
            run = first_p.add_run()
            run.text = experience_lines[0]
            run.font.name = "Arial"
            run.font.size = Pt(9)
            run.font.color.rgb = DARK_GRAY

            for line in experience_lines[1:]:
                p = exp_content_frame.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.name = "Arial"
                run.font.size = Pt(9)
                run.font.color.rgb = DARK_GRAY

        # Footer com página e copyright - diminuir para 7pt
        footer = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(5), Inches(0.3))
        footer_frame = footer.text_frame
        footer_p = footer_frame.paragraphs[0]
        footer_run = footer_p.add_run()
        footer_run.text = f"{idx+1}     | Copyright © 2022 Accenture. All rights reserved."
        footer_font = footer_run.font
        footer_font.name = "Arial"
        footer_font.size = Pt(9)
        footer_font.color.rgb = DARK_GRAY

        # Logo Accenture Technology (direita, alinhado) - diminuir para 12pt
        logo_text = slide.shapes.add_textbox(Inches(6.4), Inches(6.9), Inches(3), Inches(0.3))
        logo_frame = logo_text.text_frame
        logo_p = logo_frame.paragraphs[0]
        logo_p.alignment = PP_ALIGN.RIGHT

        logo_run = logo_p.add_run()
        logo_run.text = "Accenture"
        logo_run.font.name = "Arial"
        logo_run.font.size = Pt(12)
        logo_run.font.bold = True

        tech_run = logo_p.add_run()
        tech_run.text = "Technology"
        tech_run.font.name = "Arial"
        tech_run.font.size = Pt(12)
        tech_run.font.bold = True
        tech_run.font.color.rgb = PURPLE_COLOR

    prs.save(output_path)
    return f"✅ Presentation successfully created at: {output_path}"


# Exemplo de uso omitido (preservar seu bloco if __name__ == "__main__")


if __name__ == "__main__":
    from lerexcel import ler_dados_excel
    from gerartabela import gerar_resumo_trabalhadores
    
    try:
        # Create a simple placeholder circle image if it doesn't exist
        if not os.path.exists('placeholder_circle.png'):
            print("Note: For better results, create a 'placeholder_circle.png' file with a blue circle outline")
        
        # Load and process data
        df = ler_dados_excel("Skills.xlsx")
        resumo = gerar_resumo_trabalhadores(df)
        
        # Generate the presentation
        result = create_formatted_cv_presentation(resumo)
        print(result)
    except Exception as e:
        print(f"❌ Error: {str(e)}")