from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import os


def create_individual_presentation(employee, idx, output_folder="Employee_Presentations"):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Cores
    PURPLE_COLOR = RGBColor(128, 0, 128)
    DARK_GRAY = RGBColor(64, 64, 64)
    BLUE_COLOR = RGBColor(0, 0, 255)

    slide = prs.slides.add_slide(blank_layout)

    # Círculo azul
    left = Inches(0.4)
    top = Inches(0.4)
    width = height = Inches(1)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    circle.fill.background()
    circle.line.color.rgb = BLUE_COLOR
    circle.line.width = Pt(1.5)

    # Texto dentro do círculo
    pic_text = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(0.5))
    pic_text_frame = pic_text.text_frame
    pic_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = pic_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Place\nPicture here\nOr Delete"
    font = run.font
    font.name = "Arial"
    font.size = Pt(8)
    font.bold = True
    font.color.rgb = BLUE_COLOR

    # Nome
    name_box = slide.shapes.add_textbox(Inches(1.7), Inches(0.4), Inches(6), Inches(0.8))
    name_frame = name_box.text_frame
    name_p = name_frame.paragraphs[0]
    name_p.alignment = PP_ALIGN.LEFT
    name_run = name_p.add_run()
    name_run.text = employee["First Name Last Name"]
    name_font = name_run.font
    name_font.name = "Arial"
    name_font.size = Pt(28)
    name_font.color.rgb = DARK_GRAY

    # Cargo
    job_box = slide.shapes.add_textbox(Inches(1.7), Inches(1.1), Inches(6), Inches(0.5))
    job_frame = job_box.text_frame
    job_p = job_frame.paragraphs[0]
    job_p.alignment = PP_ALIGN.LEFT
    job_run = job_p.add_run()
    job_run.text = employee["Job Title / Role"]
    job_font = job_run.font
    job_font.name = "Arial"
    job_font.size = Pt(14)
    job_font.color.rgb = DARK_GRAY

    # Cabeçalhos + Conteúdos: Profile, Education, Skills, Experience
    def add_header(slide, text, left, top):
        box = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.4))
        frame = box.text_frame
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = "Arial"
        font.size = Pt(18)
        font.bold = True
        font.color.rgb = PURPLE_COLOR

    def add_text(slide, content, left, top, width, height, font_size=9):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = content
        font = run.font
        font.name = "Arial"
        font.size = Pt(font_size)
        font.color.rgb = DARK_GRAY

    add_header(slide, "Profile", Inches(0.4), Inches(1.8))
    add_text(slide, employee["Profile"], Inches(0.4), Inches(2.2), Inches(5), Inches(1.1))

    add_header(slide, "Education", Inches(0.4), Inches(2.8))
    add_text(slide, employee["Education"], Inches(0.4), Inches(3.2), Inches(5), Inches(1.0))

    # Skills em duas colunas
    add_header(slide, "Relevant Skills & Qualifications", Inches(0.4), Inches(3.8))
    skills_list = employee["Relevant Skills & Qualifications"].split('\n')
    max_rows = 8
    rows1 = min(len(skills_list), max_rows)
    rows2 = len(skills_list) - rows1

    table1 = slide.shapes.add_table(rows1, 1, Inches(0.5), Inches(4.2), Inches(2.2), Inches(1.5)).table
    for i in range(rows1):
        cell = table1.cell(i, 0)
        cell.text = skills_list[i]
        cell.fill.background()
        for p in cell.text_frame.paragraphs:
            p.font.name = "Arial"
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GRAY

    if rows2 > 0:
        table2 = slide.shapes.add_table(rows2, 1, Inches(3.0), Inches(4.2), Inches(2.2), Inches(1.5)).table
        for i in range(rows2):
            cell = table2.cell(i, 0)
            cell.text = skills_list[i + max_rows]
            cell.fill.background()
            for p in cell.text_frame.paragraphs:
                p.font.name = "Arial"
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_GRAY

    # Experience
    add_header(slide, "Relevant Experience", Inches(6.0), Inches(1.8))
    exp_box = slide.shapes.add_textbox(Inches(6.0), Inches(2.3), Inches(3.5), Inches(3.6))
    exp_frame = exp_box.text_frame
    exp_frame.word_wrap = True
    lines = employee["Relevant Experience"].split('\n')
    if lines:
        p = exp_frame.paragraphs[0]
        p.clear()
        r = p.add_run()
        r.text = lines[0]
        r.font.name = "Arial"
        r.font.size = Pt(9)
        r.font.color.rgb = DARK_GRAY
        for line in lines[1:]:
            p = exp_frame.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.name = "Arial"
            r.font.size = Pt(9)
            r.font.color.rgb = DARK_GRAY

    # Rodapé
    footer = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(5), Inches(0.3))
    footer_p = footer.text_frame.paragraphs[0]
    run = footer_p.add_run()
    run.text = f"{idx+1}     | Copyright © 2022 Accenture. All rights reserved."
    font = run.font
    font.name = "Arial"
    font.size = Pt(9)
    font.color.rgb = DARK_GRAY

    logo = slide.shapes.add_textbox(Inches(6.4), Inches(6.9), Inches(3), Inches(0.3))
    logo_p = logo.text_frame.paragraphs[0]
    logo_p.alignment = PP_ALIGN.RIGHT
    r1 = logo_p.add_run()
    r1.text = "Accenture"
    r1.font.name = "Arial"
    r1.font.size = Pt(12)
    r1.font.bold = True
    r2 = logo_p.add_run()
    r2.text = "Technology"
    r2.font.name = "Arial"
    r2.font.size = Pt(12)
    r2.font.bold = True
    r2.font.color.rgb = PURPLE_COLOR

    filename = f"{employee['First Name Last Name'].replace(' ', '_')}_{idx+1}.pptx"
    filepath = os.path.join(output_folder, filename)
    prs.save(filepath)

    return filepath

def create_all_presentations(data_df, output_folder="Employee_Presentations"):
    created_files = []
    for idx, employee in data_df.iterrows():
        filepath = create_individual_presentation(employee, idx, output_folder)
        created_files.append(filepath)
    return created_files

if __name__ == "__main__":
    from lerexcel import ler_dados_excel
    from gerartabela import gerar_resumo_trabalhadores

    try:
        df = ler_dados_excel("Skills.xlsx")
        resumo = gerar_resumo_trabalhadores(df)
        
        arquivos = create_all_presentations(resumo)
        print(f"✅ Created {len(arquivos)} presentations.")
        for f in arquivos:
            print(f" - {f}")
    except Exception as e:
        print(f"❌ Error: {str(e)}")

