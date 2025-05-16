from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import os


def create_individual_presentation(employee, idx, output_folder="Employee_Presentations"):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    prs = Presentation()
    # Definindo tamanho 16:9 (largura x altura em polegadas)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[6]

    # Cores
    PURPLE_COLOR = RGBColor(128, 0, 128)
    DARK_GRAY = RGBColor(64, 64, 64)
    BLUE_COLOR = RGBColor(0, 0, 255)

    slide = prs.slides.add_slide(blank_layout)

    # Círculo azul
    left = Inches(0.6)
    top = Inches(0.6)
    width = height = Inches(1.2)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    circle.fill.background()
    circle.line.color.rgb = BLUE_COLOR
    circle.line.width = Pt(1.5)

    # Texto dentro do círculo
    pic_text = slide.shapes.add_textbox(left, top + Inches(0.36), width, Inches(0.6))
    pic_text_frame = pic_text.text_frame
    pic_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = pic_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Place\nPicture here\nOr Delete"
    font = run.font
    font.name = "Arial"
    font.size = Pt(9)
    font.bold = True
    font.color.rgb = BLUE_COLOR

    # Nome
    name_box = slide.shapes.add_textbox(Inches(2.4), Inches(0.6), Inches(8), Inches(1.0))
    name_frame = name_box.text_frame
    name_p = name_frame.paragraphs[0]
    name_p.alignment = PP_ALIGN.LEFT
    name_run = name_p.add_run()
    name_run.text = employee["First Name Last Name"]
    name_font = name_run.font
    name_font.name = "Arial"
    name_font.size = Pt(30)
    name_font.color.rgb = DARK_GRAY

    # Cargo
    job_box = slide.shapes.add_textbox(Inches(2.4), Inches(1.4), Inches(8), Inches(0.6))
    job_frame = job_box.text_frame
    job_p = job_frame.paragraphs[0]
    job_p.alignment = PP_ALIGN.LEFT
    job_run = job_p.add_run()
    job_run.text = employee["Job Title / Role"]
    job_font = job_run.font
    job_font.name = "Arial"
    job_font.size = Pt(14)
    job_font.color.rgb = DARK_GRAY

    # Funções auxiliares
    def add_header(slide, text, left, top):
        box = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(0.5))
        frame = box.text_frame
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = "Arial"
        font.size = Pt(20)
        font.bold = True
        font.color.rgb = PURPLE_COLOR

    def add_text(slide, content, left, top, width, height, font_size=12):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = content
        font = run.font
        font.name = "Arial"
        font.size = Pt(font_size)
        font.color.rgb = DARK_GRAY

    # Conteúdos
    add_header(slide, "Profile", Inches(0.6), Inches(2.2))
    add_text(slide, employee["Profile"], Inches(0.6), Inches(2.8), Inches(6.5), Inches(1.3))

    add_header(slide, "Education", Inches(0.6), Inches(3.5))
    add_text(slide, employee["Education"], Inches(0.6), Inches(4.1), Inches(6.5), Inches(1.2))

    add_header(slide, "Relevant Skills & Qualifications", Inches(0.6), Inches(5.2))
    skills_list = employee["Relevant Skills & Qualifications"].split('\n')
    max_rows = 8
    rows1 = min(len(skills_list), max_rows)
    rows2 = len(skills_list) - rows1

    table1 = slide.shapes.add_table(rows1, 1, Inches(0.8), Inches(5.8), Inches(3.5), Inches(1.8)).table
    for i in range(rows1):
        cell = table1.cell(i, 0)
        cell.text = skills_list[i]
        cell.fill.background()
        for p in cell.text_frame.paragraphs:
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY

    if rows2 > 0:
        table2 = slide.shapes.add_table(rows2, 1, Inches(5.0), Inches(5.8), Inches(3.5), Inches(1.8)).table
        for i in range(rows2):
            cell = table2.cell(i, 0)
            cell.text = skills_list[i + max_rows]
            cell.fill.background()
            for p in cell.text_frame.paragraphs:
                p.font.name = "Arial"
                p.font.size = Pt(12)
                p.font.color.rgb = DARK_GRAY

    add_header(slide, "Relevant Experience", Inches(9.5), Inches(2.2))
    exp_box = slide.shapes.add_textbox(Inches(9.5), Inches(2.8), Inches(4.5), Inches(4.5))
    exp_frame = exp_box.text_frame
    exp_frame.word_wrap = True
    lines = employee["Relevant Experience"].split('\n')
    if lines:
        p = exp_frame.paragraphs[0]
        p.clear()
        r = p.add_run()
        r.text = lines[0]
        r.font.name = "Arial"
        r.font.size = Pt(12)
        r.font.color.rgb = DARK_GRAY
        for line in lines[1:]:
            p = exp_frame.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.name = "Arial"
            r.font.size = Pt(12)
            r.font.color.rgb = DARK_GRAY

    # Rodapé
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(8.5), Inches(7), Inches(0.5))
    footer_p = footer.text_frame.paragraphs[0]
    run = footer_p.add_run()
    run.text = f"{idx+1}     | Copyright © 2022 Accenture. All rights reserved."
    font = run.font
    font.name = "Arial"
    font.size = Pt(9)
    font.color.rgb = DARK_GRAY

    logo = slide.shapes.add_textbox(Inches(10.0), Inches(8.5), Inches(6), Inches(0.5))
    logo_p = logo.text_frame.paragraphs[0]
    logo_p.alignment = PP_ALIGN.RIGHT
    r1 = logo_p.add_run()
    r1.text = "Accenture"
    r1.font.name = "Arial"
    r1.font.size = Pt(20)
    r1.font.bold = True
    r2 = logo_p.add_run()
    r2.text = "Technology"
    r2.font.name = "Arial"
    r2.font.size = Pt(20)
    r2.font.bold = True
    r2.font.color.rgb = PURPLE_COLOR

    filename = f"{employee['First Name Last Name'].replace(' ', '_')}_{idx+1}.pptx"
    filepath = os.path.join(output_folder, filename)
    prs.save(filepath)

    return filepath

# Mantém create_all_presentations inalterado

def create_all_presentations(data_df, output_folder="Employee_Presentations"):
    created_files = []
    for idx, employee in data_df.iterrows():
        filepath = create_individual_presentation(employee, idx, output_folder)
        created_files.append(filepath)
    return created_files

if __name__ == "__main__":
    from lerexcel import ler_dados_excel
    from gerartabela import gerar_resumo_trabalhadores

    try:
        df = ler_dados_excel("Skills.xlsx")
        resumo = gerar_resumo_trabalhadores(df)
        
        arquivos = create_all_presentations(resumo)
        print(f"✅ Created {len(arquivos)} presentations.")
        for f in arquivos:
            print(f" - {f}")
    except Exception as e:
        print(f"❌ Error: {str(e)}")

