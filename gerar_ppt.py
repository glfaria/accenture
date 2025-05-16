import os
import pandas as pd
from collections import defaultdict
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import comtypes.client
from PyPDF2 import PdfMerger

# === CONFIGURAÇÕES ===
excel_path = "Skills.xlsx"
template_path = "CV.pptx"
output_dir = "output"
pptx_dir = os.path.join(output_dir, "pptx")
pdf_dir = os.path.join(output_dir, "pdf")
final_pdf = os.path.join(output_dir, "CVs_final.pdf")

os.makedirs(pptx_dir, exist_ok=True)
os.makedirs(pdf_dir, exist_ok=True)

# === LER EXCEL ===
df = pd.read_excel(excel_path)
grouped = defaultdict(list)
for _, row in df.iterrows():
    grouped[row["Worker Name"]].append(row)

# === FUNÇÃO PARA PREENCHER TEMPLATE ===
def preencher_cv(worker_name, rows):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    row = rows[0]
    job_title = row['Job Title']
    industry = row['Industry Networks']
    function = row['Function Networks']
    tech = row['Technology Networks']
    job_family = row['Job Family']
    profile_text = f"Professional with experience in {industry}, specialized in {function} and {tech}. Works in the {job_family} domain."

    skills_text = ""
    seen = set()
    for r in rows:
        skill = str(r['Skill'])
        proficiency = str(r['Skill Proficiency'])
        if skill not in seen:
            skills_text += f"- {skill} ({proficiency})\n"
            seen.add(skill)

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text = shape.text.strip()
            if "First Name Last Name" in text:
                shape.text = worker_name
            elif "Job Title / Role" in text:
                shape.text = job_title
            elif "Profile" in text:
                shape.text = profile_text
            elif "Relevant Skills & Qualifications" in text:
                shape.text = skills_text.strip()
            elif "Relevant Experience" in text:
                shape.text = "Experience details available upon request."
            elif "Education" in text:
                shape.text = "Education details available upon request."

    pptx_path = os.path.join(pptx_dir, f"CV_{worker_name.replace(' ', '_')}.pptx")
    prs.save(pptx_path)
    return pptx_path

# === PREENCHER TODOS OS CVs ===
pptx_paths = []
for worker, rows in grouped.items():
    pptx_path = preencher_cv(worker, rows)
    pptx_paths.append(pptx_path)

# === CONVERTER PPTX ➝ PDF (Windows + PowerPoint) ===
def pptx_to_pdf(input_path, output_path):
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    deck = powerpoint.Presentations.Open(input_path, WithWindow=False)
    deck.SaveAs(output_path, FileFormat=32)  # 32 = PDF
    deck.Close()
    powerpoint.Quit()

pdf_paths = []
for pptx_file in pptx_paths:
    filename = os.path.splitext(os.path.basename(pptx_file))[0] + ".pdf"
    pdf_path = os.path.join(pdf_dir, filename)
    pptx_to_pdf(pptx_file, pdf_path)
    pdf_paths.append(pdf_path)

# === UNIR TODOS OS PDFs ===
merger = PdfMerger()
for pdf in pdf_paths:
    merger.append(pdf)
merger.write(final_pdf)
merger.close()

print(f"✅ PDF final gerado: {final_pdf}")
