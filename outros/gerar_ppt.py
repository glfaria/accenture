import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PyPDF2 import PdfMerger



# === CONFIGURAÇÕES ===
excel_path = "Skills.xlsx"
output_dir = "output"
pptx_dir = os.path.join(output_dir, "pptx")
pdf_dir = os.path.join(output_dir, "pdf")
final_pdf = os.path.join(output_dir, "CVs_final.pdf")

os.makedirs(pptx_dir, exist_ok=True)
os.makedirs(pdf_dir, exist_ok=True)

def criar_template_cv():
    """Cria um novo template de CV do zero baseado no design da imagem"""
    prs = Presentation()
    
    # Definir tamanho do slide como A4 (proporção 16:9 mais aproximada)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Adicionar slide
    blank_slide_layout = prs.slide_layouts[6]  # Layout em branco
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Configuração de cores
    cor_roxa = RGBColor(128, 0, 128)  # Roxo para títulos
    cor_cinza = RGBColor(100, 100, 100)  # Cinza para texto normal
    cor_fundo_titulo = RGBColor(240, 240, 240)  # Cinza claro para fundo de títulos
    
    # ==== ÁREA DE CABEÇALHO ====
    
    # Placeholder para imagem (círculo)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(1)
    height = Inches(1)
    
    img_placeholder = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
    img_placeholder.line.color.rgb = RGBColor(90, 90, 255)
    
    # Adicionar texto dentro do círculo
    text_frame = img_placeholder.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Place\nPicture here\nOr Delete"
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    
    for run in p.runs:
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Nome da pessoa
    left = Inches(2)
    top = Inches(0.5)
    width = Inches(7)
    height = Inches(0.8)
    
    nome_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = nome_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "First Name Last Name"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.color.rgb = cor_cinza
    
    # Cargo
    left = Inches(2)
    top = Inches(1.2)
    width = Inches(7)
    height = Inches(0.4)
    
    cargo_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = cargo_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Job Title / Role"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(18)
        run.font.color.rgb = cor_cinza
    
    # ==== COLUNA ESQUERDA (40% largura) ====
    left_col_width = Inches(3.8)
    
    # Seção Perfil
    top_profile = Inches(1.8)
    
    perfil_titulo = slide.shapes.add_textbox(
        Inches(0.5), top_profile, left_col_width, Inches(0.4)
    )
    text_frame = perfil_titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Profile"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(20)
        run.font.color.rgb = cor_roxa
        run.font.bold = True
    
    # Conteúdo do perfil
    perfil_conteudo = slide.shapes.add_textbox(
        Inches(0.5), top_profile + Inches(0.5), left_col_width, Inches(1.5)
    )
    text_frame = perfil_conteudo.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Experience in business and technology consulting, with initial focus on a specific industry sector. Has supported projects involving process improvements, system implementations, and organizational change. Also contributed to initiatives in areas such as financial services, energy, and the public sector."
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Seção Educação
    top_education = top_profile + Inches(2.2)
    
    edu_titulo = slide.shapes.add_textbox(
        Inches(0.5), top_education, left_col_width, Inches(0.4)
    )
    text_frame = edu_titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Education"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(20)
        run.font.color.rgb = cor_roxa
        run.font.bold = True
    
    # Conteúdo da educação
    edu_conteudo = slide.shapes.add_textbox(
        Inches(0.65), top_education + Inches(0.5), left_col_width - Inches(0.15), Inches(1.2)
    )
    text_frame = edu_conteudo.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "University Name – Postgraduate Program in Course Name – 20XX–20XX"
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    p = text_frame.add_paragraph()
    p.text = "University Name – Master's Degree in Course Name – 20XX–20XX"
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    p = text_frame.add_paragraph()
    p.text = "University Name – Bachelor's Degree in Course Name – 20XX–20XX"
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Seção Skills
    top_skills = top_education + Inches(1.8)
    
    skills_titulo = slide.shapes.add_textbox(
        Inches(0.5), top_skills, left_col_width, Inches(0.4)
    )
    text_frame = skills_titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Relevant Skills & Qualifications"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(20)
        run.font.color.rgb = cor_roxa
        run.font.bold = True
    
    # Retângulo para conteúdo de skills
    skills_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 
        Inches(0.5), 
        top_skills + Inches(0.5), 
        left_col_width, 
        Inches(2.2)
    )
    skills_box.fill.solid()
    skills_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    skills_box.line.color.rgb = RGBColor(200, 200, 200)
    
    # Título Functional/Technical
    func_titulo = slide.shapes.add_textbox(
        Inches(0.6), top_skills + Inches(0.6), Inches(1.8), Inches(0.25)
    )
    text_frame = func_titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Functional/Technical:"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.bold = True
    
    # Lista de skills técnicas
    skills_list = slide.shapes.add_textbox(
        Inches(0.8), top_skills + Inches(0.85), Inches(1.7), Inches(1.8)
    )
    text_frame = skills_list.text_frame
    text_frame.word_wrap = True
    
    skills_items = [
        "Requirements Analysis",
        "User stories creation",
        "Process Modelling (BPMN)",
        "Process Optimization Methodologies",
        "JIRA/Confluence/Azure tools",
        "Project Management",
        "MS project tool",
        "Digital Transformation",
        "Agile Methodologies"
    ]
    
    for i, skill in enumerate(skills_items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = f"• {skill}"
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = cor_cinza
    
    # Título Industries
    ind_titulo = slide.shapes.add_textbox(
        Inches(2.6), top_skills + Inches(0.6), Inches(1.5), Inches(0.25)
    )
    text_frame = ind_titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Industries:"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.bold = True
    
    # Lista de indústrias
    ind_list = slide.shapes.add_textbox(
        Inches(2.8), top_skills + Inches(0.85), Inches(1.3), Inches(0.6)
    )
    text_frame = ind_list.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "• Insurance"
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    p = text_frame.add_paragraph()
    p.text = "• Public Sector"
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Caixa de línguas
    lang_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 
        Inches(2.6), 
        top_skills + Inches(1.6), 
        Inches(1.5), 
        Inches(0.9)
    )
    lang_box.fill.solid()
    lang_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    lang_box.line.color.rgb = RGBColor(200, 200, 200)
    
    # Título Languages
    lang_titulo = slide.shapes.add_textbox(
        Inches(2.7), top_skills + Inches(1.7), Inches(1.3), Inches(0.25)
    )
    text_frame = lang_titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Languages:"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.bold = True
    
    # Lista de línguas
    lang_list = slide.shapes.add_textbox(
        Inches(2.9), top_skills + Inches(1.95), Inches(1.1), Inches(0.5)
    )
    text_frame = lang_list.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "• Portuguese (Native)"
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    p = text_frame.add_paragraph()
    p.text = "• English (C1)"
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Rodapé
    footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.2), Inches(9), Inches(0.3)
    )
    text_frame = footer.text_frame
    p = text_frame.paragraphs[0]
    p.text = "1       | Copyright © 2022 Accenture. All rights reserved."
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(8)
        run.font.color.rgb = cor_cinza
    
    # Logotipo Accenture Technology
    logo = slide.shapes.add_textbox(
        Inches(7), Inches(7.15), Inches(2.5), Inches(0.35)
    )
    text_frame = logo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Accenture"
    p.alignment = PP_ALIGN.RIGHT
    
    run = p.runs[0]
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0, 0, 0)
    run.font.bold = True
    
    p.add_run()
    run = p.runs[1]
    run.text = "Technology"
    run.font.size = Pt(16)
    run.font.color.rgb = cor_roxa
    run.font.bold = True
    
    # ==== COLUNA DIREITA (60% largura) ====
    right_col_left = Inches(4.6)
    right_col_width = Inches(4.9)
    
    # Título Experiência
    exp_titulo = slide.shapes.add_textbox(
        right_col_left, Inches(1.8), right_col_width, Inches(0.4)
    )
    text_frame = exp_titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Relevant Experience"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(20)
        run.font.color.rgb = cor_roxa
        run.font.bold = True
    
    # Experiência 1
    exp_top = Inches(2.3)
    
    # Data
    exp1_data = slide.shapes.add_textbox(
        right_col_left, exp_top, right_col_width, Inches(0.2)
    )
    text_frame = exp1_data.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Jun 2024 - Sep 2024"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Cargo
    exp1_cargo = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.2), right_col_width, Inches(0.2)
    )
    text_frame = exp1_cargo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Business Analyst"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
        run.font.bold = True
    
    # Descrição
    exp1_desc = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.4), right_col_width, Inches(0.4)
    )
    text_frame = exp1_desc.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Joined a business analysis team on a project focused on creating an innovative platform to support the mapping, interpretation, and strategic valuation of territorial data."
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Experiência 2
    exp_top = exp_top + Inches(0.9)
    
    # Data
    exp2_data = slide.shapes.add_textbox(
        right_col_left, exp_top, right_col_width, Inches(0.2)
    )
    text_frame = exp2_data.text_frame
    p = text_frame.paragraphs[0]
    p.text = "May 2023 - Dec 2023"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Cargo
    exp2_cargo = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.2), right_col_width, Inches(0.2)
    )
    text_frame = exp2_cargo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Product Owner"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
        run.font.bold = True
    
    # Descrição
    exp2_desc = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.4), right_col_width, Inches(0.4)
    )
    text_frame = exp2_desc.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Worked alongside a product team, interacting with stakeholders to define priorities and write detailed user stories for the development process."
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Experiência 3
    exp_top = exp_top + Inches(0.9)
    
    # Data
    exp3_data = slide.shapes.add_textbox(
        right_col_left, exp_top, right_col_width, Inches(0.2)
    )
    text_frame = exp3_data.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Jan 2023 - May 2023"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Cargo
    exp3_cargo = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.2), right_col_width, Inches(0.2)
    )
    text_frame = exp3_cargo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Product Owner"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
        run.font.bold = True
    
    # Descrição
    exp3_desc = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.4), right_col_width, Inches(0.4)
    )
    text_frame = exp3_desc.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Contributed to the definition of the product roadmap, aligning business priorities and stakeholder expectations to ensure value delivery."
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Experiência 4
    exp_top = exp_top + Inches(0.9)
    
    # Data
    exp4_data = slide.shapes.add_textbox(
        right_col_left, exp_top, right_col_width, Inches(0.2)
    )
    text_frame = exp4_data.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Nov 2021 - Apr 2022"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Cargo
    exp4_cargo = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.2), right_col_width, Inches(0.2)
    )
    text_frame = exp4_cargo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Senior Business Analyst"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
        run.font.bold = True
    
    # Descrição
    exp4_desc = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.4), right_col_width, Inches(0.4)
    )
    text_frame = exp4_desc.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Identified business requirements and designed solutions to support the implementation of a system for managing operational processes."
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Experiência 5
    exp_top = exp_top + Inches(0.9)
    
    # Data
    exp5_data = slide.shapes.add_textbox(
        right_col_left, exp_top, right_col_width, Inches(0.2)
    )
    text_frame = exp5_data.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Sep 2019 - Sep 2021"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Cargo
    exp5_cargo = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.2), right_col_width, Inches(0.2)
    )
    text_frame = exp5_cargo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Senior Business Analyst"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
        run.font.bold = True
    
    # Projeto 1
    exp5_proj1 = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.4), right_col_width, Inches(0.2)
    )
    text_frame = exp5_proj1.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Project 1 - Digital Transformation"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
        run.font.italic = True
    
    # Descrição Projeto 1
    exp5_proj1_desc = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(0.6), right_col_width, Inches(0.4)
    )
    text_frame = exp5_proj1_desc.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Led the creation and implementation of a process framework, including governance structures, responsibility matrices, and performance indicators. Developed a continuous improvement model to enhance organizational efficiency."
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    # Projeto 2
    exp5_proj2 = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(1.0), right_col_width, Inches(0.2)
    )
    text_frame = exp5_proj2.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Project 2 - Business Projects Management"
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
        run.font.italic = True
    
    # Descrição Projeto 2
    exp5_proj2_desc = slide.shapes.add_textbox(
        right_col_left, exp_top + Inches(1.2), right_col_width, Inches(0.7)
    )
    text_frame = exp5_proj2_desc.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Prepared business cases for new strategic initiatives. Designed and implemented operational models and business processes to support service delivery."
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    p = text_frame.add_paragraph()
    p.text = "Supported the definition of system requirements, contributed to integration and user acceptance testing, and participated in system assessment activities for new solution implementations."
    p.alignment = PP_ALIGN.LEFT
    
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = cor_cinza
    
    return prs

def ler_excel_e_preencher_cvs():
    """Lê dados do Excel e preenche os CVs com os dados apropriados"""
    try:
        df = pd.read_excel(excel_path)
        df.columns = df.columns.str.strip()  # remove espaços extras
        
        # Agrupar por nome do trabalhador
        grouped = df.groupby("Worker Name")
        
        pptx_paths = []
        
        for worker_name, group in grouped:
            # Criar um novo template para cada trabalhador
            prs = criar_template_cv()
            slide = prs.slides[0]
            
            # Extrair informações do primeiro registro do trabalhador
            first_row = group.iloc[0]
            
            job_title = first_row.get('Job Title', '')
            
            # Coletar skills do trabalhador
            skills = []
            for _, row in group.iterrows():
                skill = str(row.get('Skill', ''))
                proficiency = str(row.get('Skill Proficiency', ''))
                if skill and skill not in [s[0] for s in skills]:
                    skills.append((skill, proficiency))
            
            # Função para encontrar e atualizar uma forma baseada no seu texto
            def update_shape_text(text_to_find, new_text):
                for shape in slide.shapes:
                    if shape.has_text_frame and text_to_find in shape.text:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if text_to_find in run.text:
                                    run.text = new_text
                                    return True
                return False
            
            # Atualizar nome e cargo
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "First Name Last Name" in shape.text:
                        text_frame = shape.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        p.text = worker_name
                        
                        for run in p.runs:
                            run.font.size = Pt(36)
                            run.font.color.rgb = RGBColor(100, 100, 100)
                    
                    elif "Job Title / Role" in shape.text:
                        text_frame = shape.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        p.text = job_title
                        
                        for run in p.runs:
                            run.font.size = Pt(18)
                            run.font.color.rgb = RGBColor(100, 100, 100)
            
            # Construir perfil baseado nos dados
            profile_text = f"Professional with experience in {first_row.get('Industry Networks', '')} industry, "
            profile_text += f"specialized in {first_row.get('Function Networks', '')} and {first_row.get('Technology Networks', '')}. "
            profile_text += f"Works in the {first_row.get('Job Family', '')} domain with focus on business analysis and process improvement."
            
            # Atualizar texto do perfil
            for shape in slide.shapes:
                if shape.has_text_frame and "Profile" in shape.text and len(shape.text) < 10:
                    # Encontrou o título do perfil, agora procurar o conteúdo abaixo
                    shape_left = shape.left
                    shape_top = shape.top
                    shape_width = shape.width
                    
                    # Procurar a caixa de texto do conteúdo do perfil
                    for content_shape in slide.shapes:
                        if (content_shape.has_text_frame and 
                            content_shape.left == shape_left and 
                            content_shape.top > shape_top and 
                            content_shape.top < shape_top + Inches(1) and
                            len(content_shape.text) > 20):
                            
                            text_frame = content_shape.text_frame
                            text_frame.clear()
                            p = text_frame.paragraphs[0]
                            p.text = profile_text
                            
                            for run in p.runs:
                                run.font.size = Pt(10)
                                run.font.color.rgb = RGBColor(100, 100, 100)
                            
                            break
            
            # Atualizar habilidades
            skill_items = []
            for skill, proficiency in skills[:9]:  # Limitar a 9 habilidades
                skill_items.append(f"• {skill} ({proficiency})")
            
            # Encontrar a caixa de habilidades e atualizar
            for shape in slide.shapes:
                if (shape.has_text_frame and 
                    "Requirements Analysis" in shape.text and
                    "User stories" in shape.text):
                    
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    for i, skill_item in enumerate(skill_items):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = skill_item
                        p.alignment = PP_ALIGN.LEFT
                        
                        for run in p.runs:
                            run.font.size = Pt(10)
                            run.font.color.rgb = RGBColor(100, 100, 100)
                    
                    # Se tiver menos de 9 habilidades, preencher com espaços
                    for i in range(len(skill_items), 9):
                        p = text_frame.add_paragraph()
                        p.text = " "
                        for run in p.runs:
                            run.font.size = Pt(10)
            
            # Salvar o arquivo PPTX
            output_filename = f"CV_{worker_name.replace(' ', '_')}.pptx"
            output_path = os.path.join(pptx_dir, output_filename)
            prs.save(output_path)
            pptx_paths.append(output_path)
            
            print(f"CV criado para {worker_name}: {output_path}")
        
        return pptx_paths
    
    except Exception as e:
        print(f"Erro ao processar Excel: {e}")
        return []

def pptx_to_pdf(input_path, output_dir):
    """Converte PPTX para PDF usando LibreOffice"""
    try:
        filename = os.path.basename(input_path)
        name_without_ext = os.path.splitext(filename)[0]
        output_path = os.path.join(output_dir, f"{name_without_ext}.pdf")
        
        # Comando para LibreOffice (Linux/Mac)
        cmd = f'libreoffice --headless --convert-to pdf --outdir "{output_dir}" "{input_path}"'
        
        # Se estiver no Windows, ajustar o caminho do libreoffice
        import platform
        if platform.system() == "Windows":
            # Ajuste este caminho para a sua instalação do LibreOffice no Windows
            libreoffice_path = r'C:\Program Files\LibreOffice\program\soffice.exe'
            cmd = f'"{libreoffice_path}" --headless --convert-to pdf --outdir "{output_dir}" "{input_path}"'
        
        os.system(cmd)
        
        # Verificar se o PDF foi criado
        if os.path.exists(output_path):
            return output_path
        else:
            print(f"Falha ao converter {input_path} para PDF")
            return None
    
    except Exception as e:
        print(f"Erro ao converter para PDF: {e}")
        return None

def unir_pdfs(pdf_paths, output_path):
    """Une vários PDFs em um único arquivo"""
    try:
        merger = PdfMerger()
        
        for pdf in pdf_paths:
            if pdf and os.path.exists(pdf):
                merger.append(pdf)
            else:
                print(f"PDF não encontrado: {pdf}")
        
        if len(merger.pages) > 0:
            merger.write(output_path)
            merger.close()
            print(f"PDFs unidos com sucesso: {output_path}")
            return True
        else:
            print("Nenhum PDF válido para unir")
            return False
    
    except Exception as e:
        print(f"Erro ao unir PDFs: {e}")
        return False

def main():
    """Função principal do programa"""
    # Criar diretórios de saída
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs(pptx_dir, exist_ok=True)
    os.makedirs(pdf_dir, exist_ok=True)
    
    print("\n=== Gerador de CVs ===\n")
    
    # Gerar um template de exemplo
    print("Gerando template de exemplo...")
    template = criar_template_cv()
    template_path = os.path.join(output_dir, "CV_Template.pptx")
    template.save(template_path)
    print(f"Template salvo em: {template_path}")
    
    # Verificar se o arquivo Excel existe
    if not os.path.exists(excel_path):
        print(f"Arquivo Excel não encontrado: {excel_path}")
        print("Geração apenas do template concluída.")
        return
    
    # Ler Excel e gerar CVs
    print("\nProcessando dados do Excel e gerando CVs...")
    pptx_paths = ler_excel_e_preencher_cvs()
    
    if not pptx_paths:
        print("Nenhum CV foi gerado. Verifique o arquivo Excel.")
        return
    
    # Converter para PDF (opcional)
    pdf_paths = []
    print("\nConvertendo para PDF...")
    for pptx_path in pptx_paths:
        pdf_path = pptx_to_pdf(pptx_path, pdf_dir)
        if pdf_path:
            pdf_paths.append(pdf_path)
    
    # Unir PDFs (opcional)
    if pdf_paths:
        print("\nUnindo PDFs...")
        if unir_pdfs(pdf_paths, final_pdf):
            print(f"\nProcesso concluído com sucesso! PDF final: {final_pdf}")
        else:
            print("\nFalha ao unir PDFs.")
    else:
        print("\nNenhum PDF gerado para unir.")
    
    print("\nProcesso finalizado!")

if __name__ == "__main__":
    main()